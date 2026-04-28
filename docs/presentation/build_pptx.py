"""
Generate MoTitle feature presentation PPTX from captured screenshots.

Run:
    python3 docs/presentation/build_pptx.py
Output:
    docs/presentation/MoTitle_Feature_Presentation.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

REPO = Path(__file__).resolve().parent.parent.parent
SHOTS = REPO / "docs/presentation/screenshots"
OUT = REPO / "docs/presentation/MoTitle_Feature_Presentation.pptx"

# ---------- Theme ----------
BG = RGBColor(0x10, 0x10, 0x18)
PANEL = RGBColor(0x18, 0x18, 0x22)
ACCENT = RGBColor(0x6E, 0x5B, 0xE0)        # purple
ACCENT_2 = RGBColor(0xF9, 0xE2, 0xAF)      # amber
SUCCESS = RGBColor(0x4A, 0xDE, 0x80)       # green
WHITE = RGBColor(0xF5, 0xF5, 0xFA)
DIM = RGBColor(0xA0, 0xA0, 0xB0)
DIVIDER = RGBColor(0x38, 0x38, 0x46)

# 16:9 — 13.333" x 7.5"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_text(slide, x, y, w, h, text, *,
             size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             font="Helvetica Neue"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return tb


def add_accent_bar(slide, x, y, w, h, color=ACCENT):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def add_chip(slide, x, y, text, color=ACCENT, fg=WHITE, size=11):
    """Pill / chip / tag."""
    width = Inches(0.05 + len(text) * 0.10)
    h = Inches(0.32)
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, h)
    chip.fill.solid()
    chip.fill.fore_color.rgb = color
    chip.line.fill.background()
    chip.shadow.inherit = False
    chip.adjustments[0] = 0.5
    tf = chip.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = fg
    return chip


def add_image(slide, path, x, y, w=None, h=None):
    if not Path(path).exists():
        return None
    if w and h:
        return slide.shapes.add_picture(str(path), x, y, width=w, height=h)
    elif w:
        return slide.shapes.add_picture(str(path), x, y, width=w)
    elif h:
        return slide.shapes.add_picture(str(path), x, y, height=h)
    else:
        return slide.shapes.add_picture(str(path), x, y)


def add_header(slide, eyebrow, title, subtitle=None):
    add_accent_bar(slide, Inches(0.5), Inches(0.5), Inches(0.08), Inches(0.5))
    add_text(slide, Inches(0.7), Inches(0.45), Inches(8), Inches(0.4),
             eyebrow, size=12, bold=True, color=ACCENT_2)
    add_text(slide, Inches(0.7), Inches(0.78), Inches(12), Inches(0.7),
             title, size=30, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.7), Inches(1.42), Inches(12), Inches(0.4),
                 subtitle, size=14, color=DIM)


def add_footer(slide, page_no, total):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
             "MoTitle — 廣播字幕製作 Pipeline", size=10, color=DIM)
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.6), Inches(0.3),
             f"{page_no} / {total}", size=10, color=DIM, align=PP_ALIGN.RIGHT)


# ============================================================
# Slide builders
# ============================================================
SLIDES = []


def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    # Diagonal accent panel
    panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.6), SLIDE_W, Inches(2.5))
    panel.fill.solid()
    panel.fill.fore_color.rgb = PANEL
    panel.line.fill.background()
    panel.shadow.inherit = False
    add_accent_bar(s, Inches(0.6), Inches(2.85), Inches(0.12), Inches(2.0))
    add_text(s, Inches(0.9), Inches(2.85), Inches(11.5), Inches(0.5),
             "Broadcast-Grade Subtitle Pipeline", size=16, bold=True, color=ACCENT_2)
    add_text(s, Inches(0.9), Inches(3.35), Inches(12), Inches(1.4),
             "MoTitle", size=72, bold=True, color=WHITE)
    add_text(s, Inches(0.9), Inches(4.55), Inches(12), Inches(0.5),
             "本地優先 · 模塊化 Pipeline · 廣播交付級輸出", size=18, color=DIM)
    add_chip(s, Inches(0.9), Inches(5.4), " 100% 本地可運行 ", color=SUCCESS, fg=BG)
    add_chip(s, Inches(2.5), Inches(5.4), " 零月費 ", color=ACCENT, fg=WHITE)
    add_chip(s, Inches(3.65), Inches(5.4), " 廣播交付級 ", color=ACCENT_2, fg=BG)
    add_chip(s, Inches(5.0), Inches(5.4), " 開源 / 自託管 ", color=PANEL, fg=WHITE)
    add_text(s, Inches(0.9), Inches(6.4), Inches(12), Inches(0.4),
             "對比市面 SaaS 字幕軟件嘅完整功能介紹", size=13, color=DIM)


def slide_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "問題陳述", "市面 SaaS 字幕軟件嘅三大痛點",
               "Otter.ai · Rev · Descript · Sonix · Trint · Happy Scribe")
    items = [
        ("月費鎖定 + 按分鐘加價",
         "$20–60/月只係入場費。每分鐘音視頻再 $0.10–0.50。\n專業字幕團隊每月成本快速堆到四位數。"),
        ("數據主權失控",
         "影片必定 upload 上 cloud。受 ToS 約束、可能用作訓練、\n機密 / 法律 / 醫療項目幾乎無法用。"),
        ("黑盒模型 + Vendor lock-in",
         "ASR 模型同翻譯模型廠商鎖死，質素由佢哋決定。\n資料、流程、字幕格式都被困住，難以遷移。"),
    ]
    y = Inches(2.4)
    for i, (title, body) in enumerate(items):
        x = Inches(0.7) + Inches(4.2) * i
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x, y, Inches(4.0), Inches(3.6))
        card.fill.solid()
        card.fill.fore_color.rgb = PANEL
        card.line.fill.background()
        card.shadow.inherit = False
        card.adjustments[0] = 0.04
        # number circle
        num = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  x + Inches(0.3), y + Inches(0.3),
                                  Inches(0.55), Inches(0.55))
        num.fill.solid()
        num.fill.fore_color.rgb = ACCENT
        num.line.fill.background()
        num.shadow.inherit = False
        tf = num.text_frame
        tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(i + 1); r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = WHITE
        add_text(s, x + Inches(0.3), y + Inches(1.0), Inches(3.5), Inches(0.7),
                 title, size=16, bold=True, color=ACCENT_2)
        add_text(s, x + Inches(0.3), y + Inches(1.8), Inches(3.5), Inches(2.0),
                 body, size=12, color=WHITE)


def slide_solution(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "解決方案", "MoTitle — 完全本地、模塊化、廣播交付級")
    add_text(s, Inches(0.7), Inches(2.0), Inches(12), Inches(0.5),
             "三個核心理念：", size=16, bold=True, color=ACCENT_2)
    pillars = [
        ("本地優先",
         "ASR 用 mlx-whisper 喺 Apple Silicon Metal GPU 跑，\n翻譯用 Ollama Qwen 喺本機跑。\n影片永遠唔離開你部機。"),
        ("模塊化 Pipeline",
         "ASR、翻譯、字型、Render 四個模塊獨立配置。\nProfile 唔係預設、係樂高 — 自由組合。\n換引擎唔需要重做工作流。"),
        ("廣播交付級",
         "WYSIWYG 預覽 = 燒入結果（fontsdir + SVG paint-order）。\nMXF ProRes 422 HQ + XDCAM HD 422 直出。\n結構化 QA flags、可逆嘅詞彙表生命週期。"),
    ]
    y = Inches(2.7)
    for i, (title, body) in enumerate(pillars):
        x = Inches(0.7) + Inches(4.2) * i
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(4.0), Inches(0.06))
        bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
        bar.shadow.inherit = False
        add_text(s, x, y + Inches(0.25), Inches(4.0), Inches(0.6),
                 title, size=22, bold=True, color=WHITE)
        add_text(s, x, y + Inches(1.0), Inches(4.0), Inches(2.6),
                 body, size=13, color=DIM)


def slide_pipeline(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "Pipeline 架構", "從原始視頻到燒入字幕嘅 5 個 stage")
    stages = [
        ("📹", "原始視頻", "MP4 / MXF / MOV / MKV"),
        ("🎙", "ASR 轉錄", "Whisper · mlx-whisper · faster-whisper\nlarge-v3 / medium / small"),
        ("🌐", "翻譯", "Ollama 本地 Qwen3.5\nOpenRouter Claude / GPT / Gemini"),
        ("📝", "校對 & 詞彙表", "Proofread editor · Glossary apply\nLLM 智能替換 + 可逆 baseline"),
        ("🎬", "渲染輸出", "MP4 H.264 (CRF/CBR/2pass)\nMXF ProRes / XDCAM HD 422"),
    ]
    y = Inches(2.6)
    box_w = Inches(2.35)
    box_h = Inches(2.6)
    gap = Inches(0.18)
    for i, (icon, title, body) in enumerate(stages):
        x = Inches(0.7) + (box_w + gap) * i
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        card.fill.solid(); card.fill.fore_color.rgb = PANEL
        card.line.fill.background(); card.shadow.inherit = False
        card.adjustments[0] = 0.06
        add_text(s, x, y + Inches(0.2), box_w, Inches(0.7),
                 icon, size=32, color=ACCENT_2, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(0.95), box_w, Inches(0.5),
                 title, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.15), y + Inches(1.4), box_w - Inches(0.3), Inches(1.2),
                 body, size=10, color=DIM, align=PP_ALIGN.CENTER)
        # arrow between
        if i < len(stages) - 1:
            arrow_x = x + box_w + Inches(0.0)
            arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                        arrow_x, y + Inches(1.15), gap, Inches(0.3))
            arrow.fill.solid(); arrow.fill.fore_color.rgb = ACCENT
            arrow.line.fill.background(); arrow.shadow.inherit = False
    add_text(s, Inches(0.7), Inches(5.5), Inches(12), Inches(0.5),
             "🔁 每一 stage 嘅輸出 = 下一 stage 嘅輸入；中間任何一步可以停低 / 重做 / 換引擎",
             size=13, color=ACCENT_2, align=PP_ALIGN.CENTER)


# ---------- feature slides with screenshots ----------

def feature_slide(prs, eyebrow, title, screenshot, bullets, *, subtitle=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, eyebrow, title, subtitle)
    # Image — left half
    img_x = Inches(0.6)
    img_y = Inches(2.0)
    img_w = Inches(7.6)
    img_h = Inches(4.85)
    # Image shadow / panel
    pad = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              img_x - Inches(0.05), img_y - Inches(0.05),
                              img_w + Inches(0.1), img_h + Inches(0.1))
    pad.fill.solid(); pad.fill.fore_color.rgb = PANEL
    pad.line.color.rgb = DIVIDER; pad.line.width = Pt(1)
    pad.shadow.inherit = False
    add_image(s, screenshot, img_x, img_y, w=img_w, h=img_h)

    # Bullets — right column
    bx = Inches(8.5)
    by = Inches(2.0)
    bw = Inches(4.4)
    add_text(s, bx, by, bw, Inches(0.5), "重點功能", size=14, bold=True, color=ACCENT_2)
    cy = by + Inches(0.55)
    for label, body in bullets:
        # bullet pill
        add_chip(s, bx, cy, f" {label} ", color=ACCENT, fg=WHITE, size=10)
        add_text(s, bx, cy + Inches(0.42), bw, Inches(0.85),
                 body, size=11, color=WHITE)
        cy += Inches(1.25)


def slide_dashboard(prs):
    feature_slide(prs,
        "功能 1", "Dashboard — 一站式工作台",
        SHOTS / "01_dashboard_overview.png",
        [
            ("檔案管理", "拖放上傳、即時 transcription progress、ETA 估算、segment-level live preview。"),
            ("Profile 切換", "頂部 dropdown 即時切換 active Profile，所有後續工作流跟住換。"),
            ("狀態追蹤", "ASR 狀態、翻譯狀態、approval 進度、render 進度全部可見。"),
            ("快速渲染", "MP4 / MXF / XDCAM 三個常用 preset，一鍵燒入。"),
        ])


def slide_profile_editor(prs):
    feature_slide(prs,
        "功能 2", "Profile 編輯器 — 模塊化參數面板",
        SHOTS / "02_profile_editor.png",
        [
            ("分區編輯", "ASR、翻譯、字型、Render 四個 collapsible 區塊。"),
            ("動態 Schema", "切換引擎時參數面板按 schema 即時 re-render。"),
            ("Active 保護", "Active Profile 唔可以被刪除（防呆）。"),
            ("即時生效", "儲存後 Socket.IO 即時通知所有開緊嘅 tab。"),
        ])


def slide_proofread_overview(prs):
    feature_slide(prs,
        "功能 3", "Proofread Editor — 校對工作台",
        SHOTS / "03_proofread_overview.png",
        [
            ("並排佈局", "左視頻 + 右段落表，影片時間軸同 segment 即時 sync。"),
            ("逐句編輯", "Inline 編輯 zh_text，自動 mark approved。"),
            ("批量批核", "Approve all 一次過。可被同事檢視個體狀態。"),
            ("快捷鍵", "Cmd+F 搜尋取代、Enter 批核、Esc 取消。"),
        ])


def slide_proofread_segment(prs):
    feature_slide(prs,
        "功能 4", "Segment Detail — 句子級操作",
        SHOTS / "04_proofread_segment_detail.png",
        [
            ("英中對照", "EN 原文 + ZH 譯文 同顯，即時對齊。"),
            ("時間碼編輯", "微調 in/out 點，影片自動 seek。"),
            ("QA Flags", "[long] / [review] 結構化 flag — 唔污染譯文。"),
            ("Speaker 標籤", "可選 speaker name 標記，輔助多人對話。"),
        ])


def slide_glossary_panel(prs):
    feature_slide(prs,
        "功能 5", "詞彙表 Panel — 即時編輯 + 自動選擇",
        SHOTS / "05_glossary_panel.png",
        [
            ("Pipeline Auto-select", "Page load 時自動選中 Profile 配置嘅 glossary，省掉 dropdown 操作。"),
            ("Inline 增 / 改 / 刪", "新增、編輯、刪除全部喺 panel 入面完成，無 modal 跳轉。"),
            ("儲存 / 取消按鈕", "明顯 label 嘅 button + Enter / Esc 鍵盤支援，唔再易誤操作。"),
            ("CSV Import / Export", "舊有詞彙表可一鍵搬入；新建可一鍵備份。"),
        ])


def slide_apply_modal(prs):
    feature_slide(prs,
        "★ 旗艦功能", "Glossary Apply — EN 高亮 + LLM 智能替換",
        SHOTS / "07_glossary_apply_modal.png",
        [
            ("3 行 Row 佈局", "term header + EN 原文 + ZH 譯文。EN 永遠 highlight 詞彙；ZH 已含目標詞先 highlight。"),
            ("提示文字", "違規行 ⚠ LLM 將判斷修改位置；已符合行 ✓ 已含目標詞（強制重套切 ⚠ 強制重新套用）。"),
            ("LLM 智能替換", "唔係 string find-replace — 整段 segment 餵 LLM，保留語法時態，只改指定詞。"),
            ("Word-boundary Match", "「US」唔會 match 「must」/「trust」/「USA」（前後端 regex 完全 mirror）。"),
        ])


def slide_subtitle_settings(prs):
    feature_slide(prs,
        "功能 7", "字幕設定 Panel — Profile font 即時編輯",
        SHOTS / "08_subtitle_settings_panel.png",
        [
            ("即時 PATCH", "改字型 / 大小 / 顏色 / 輪廓 / 邊距 → 500ms debounce → 自動寫入 Profile。"),
            ("Socket.IO 跨 Tab Sync", "一個 tab 改字型、所有 tab 嘅預覽同時更新。"),
            ("Bundled Font", "Drop TTF 入 backend/assets/fonts/，前端 @font-face + libass fontsdir 同時用。"),
            ("WYSIWYG", "Browser preview 嘅每個 pixel 同最終燒入 video 完全 1:1。"),
        ])


def slide_render_mp4(prs):
    feature_slide(prs,
        "功能 8", "MP4 Render — CRF / CBR / 2-pass 三 mode",
        SHOTS / "09_render_modal_mp4.png",
        [
            ("CRF mode (預設)", "0–51 quality slider。串流 / YouTube 通用。"),
            ("CBR mode", "10–100 Mbps 嚴格定碼率。OTT 串流標準。"),
            ("2-pass mode", "兩次編碼最佳碼率分配。Master archive / 廣播交付。"),
            ("細緻控制", "Pixel format (yuv420/422/444)、H.264 profile (high/high422/high444)、level (3.1–5.2 / auto)。"),
        ])


def slide_render_mxf(prs):
    feature_slide(prs,
        "功能 9", "MXF ProRes — 廣播交付 master",
        SHOTS / "10_render_modal_mxf.png",
        [
            ("ProRes 6 個 profile", "Proxy / LT / Standard / HQ / 4444 / 4444 XQ — 對應碼率明確顯示。"),
            ("PCM 音訊", "16 / 24 / 32-bit PCM — 廣播 master 標準。"),
            ("解像度匹配", "保留原解像度 / 720p / 1080p / 4K 任揀。"),
            ("一鍵交付", "輸出檔名自動加 _subtitled，省得手動 rename。"),
        ])


def slide_render_xdcam(prs):
    feature_slide(prs,
        "功能 10", "MXF XDCAM HD 422 — Sony 廣播標準",
        SHOTS / "11_render_modal_xdcam.png",
        [
            ("Long-GOP MPEG-2 422", "MPEG-2 4:2:2 long-GOP，喺 MXF 容器，符合 SMPTE-356M。"),
            ("CBR 10–100 Mbps", "預設 50 Mbps（XDCAM HD 422 標準），可上調 100 Mbps 做 master。"),
            ("Sony 工作流相容", "可直接 import 入 Sony XDC stations、Avid、Premiere、FCP。"),
            ("全程 422 chroma", "唔係 8-bit 4:2:0 — 真正廣播 chain 友好。"),
        ])


def slide_render_2pass(prs):
    feature_slide(prs,
        "功能 11", "2-pass Encoding — 最佳碼率分配",
        SHOTS / "13_render_modal_mp4_2pass.png",
        [
            ("Pass 1 → Pass 2", "第一 pass 分析、第二 pass 用碼率 budget 智能分配難易場景。"),
            ("Concurrent-safe", "每次 render 用 unique passlogfile prefix（PID + random），唔會撞 stats file。"),
            ("自動 cleanup", ".log + .log.mbtree temp file 喺 finally block 必清。"),
            ("質素 / size 平衡", "同 CRF 同 size 比，但場景複雜處 bitrate 自動加；空鏡省 bitrate。"),
        ])


def slide_subtitle_overlay(prs):
    feature_slide(prs,
        "功能 12", "WYSIWYG 字幕預覽 — preview = burn-in",
        SHOTS / "14_subtitle_overlay.png",
        [
            ("SVG paint-order", "stroke fill linejoin=round 對齊 libass FT_Stroker 嘅幾何。"),
            ("Hardcoded viewBox 1920×1080", "= libass 嘅 PlayResX/Y。size / outline / margin pixel-perfect 1:1。"),
            ("Fontsdir 同步", "Browser @font-face 同 FFmpeg fontsdir= 用同一隻 TTF。"),
            ("結果", "Preview 同最終 H.264 / MPEG-2 4:2:0 燒入嘅每個 pixel、輪廓、字距完全一致。"),
        ])


def slide_revert(prs):
    """No screenshot — diagram-only slide for the most differentiated mechanism."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "★ 獨有機制", "Per-segment Baseline + 自動回滾",
               "業界冇商用軟件咁做 — fire-and-forget 嘅死局，喺 MoTitle 唔存在")
    # Diagram: 3-step flow
    box_w = Inches(3.6); box_h = Inches(3.4); gap = Inches(0.5)
    y = Inches(2.4)
    steps = [
        ("初次翻譯", "EN: \"in the US\"\nZH: 「喺美國」\n\nbaseline_zh = 喺美國\napplied_terms = []",
         "ZH 譯文寫入時，同步\n建立 baseline_zh 副本。"),
        ("套用詞彙表", "EN: 「the US → 美國人」\n→ LLM 智能替換\nZH: 「成為美國人」\n\napplied_terms = [(US, 美國人)]",
         "zh_text 更新；baseline_zh\n保持唔變；applied_terms\n紀錄套用嘅 (term_en, term_zh)。"),
        ("詞彙表改動 → 自動回滾", "從 glossary 刪除 \"US → 美國人\"\n下次 scan 時 lazy revert：\nzh_text ← baseline_zh\n→ 「喺美國」",
         "唔需要 manual undo。\n手動編輯會 update baseline，\n變返新嘅源頭。"),
    ]
    for i, (title, body, hint) in enumerate(steps):
        x = Inches(0.7) + (box_w + gap) * i
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        card.fill.solid(); card.fill.fore_color.rgb = PANEL
        card.line.fill.background(); card.shadow.inherit = False
        card.adjustments[0] = 0.04
        add_chip(s, x + Inches(0.3), y + Inches(0.25),
                 f" Step {i+1} ", color=ACCENT, fg=WHITE, size=10)
        add_text(s, x + Inches(0.3), y + Inches(0.7), box_w - Inches(0.6), Inches(0.5),
                 title, size=15, bold=True, color=ACCENT_2)
        add_text(s, x + Inches(0.3), y + Inches(1.25), box_w - Inches(0.6), Inches(1.4),
                 body, size=10, color=WHITE, font="Menlo")
        add_text(s, x + Inches(0.3), y + Inches(2.55), box_w - Inches(0.6), Inches(0.7),
                 hint, size=10, color=DIM)
        # arrow between
        if i < len(steps) - 1:
            arrow_x = x + box_w + Inches(0.05)
            arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                        arrow_x, y + Inches(1.55), gap - Inches(0.1), Inches(0.3))
            arrow.fill.solid(); arrow.fill.fore_color.rgb = ACCENT
            arrow.line.fill.background(); arrow.shadow.inherit = False


def slide_comparison(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "對比", "MoTitle vs SaaS 字幕軟件",
               "Otter.ai · Rev · Descript · Sonix · Trint · Happy Scribe …")
    rows = [
        ("維度", "市場 SaaS", "MoTitle"),
        ("收費模式", "$20–60/月 + 按分鐘", "開源免費"),
        ("數據去向", "影片上 cloud", "永遠喺本機"),
        ("ASR 模型", "廠商鎖死（黑盒）", "tiny / medium / large-v3 自選"),
        ("翻譯模型", "單一 provider", "Ollama 本地 + OpenRouter 9 個 model"),
        ("廣播輸出", "通常只出 MP4", "MXF ProRes + XDCAM HD 422"),
        ("Preview fidelity", "預覽 ≠ 最終輸出", "preview = burn-in（pixel-perfect）"),
        ("字幕風格 prompt", "黑盒", "繁體廣播 few-shot prompt 開放可改"),
        ("Air-gap mode", "唔可行", "100% offline 可運行"),
        ("Vendor lock-in", "高", "零（plain JSON / 開放標準）"),
    ]
    # Table
    rows_count = len(rows)
    cols_count = 3
    tbl_x = Inches(0.6); tbl_y = Inches(2.0)
    tbl_w = Inches(12.1); tbl_h = Inches(4.7)
    tbl = s.shapes.add_table(rows_count, cols_count, tbl_x, tbl_y, tbl_w, tbl_h).table
    tbl.columns[0].width = Inches(2.6)
    tbl.columns[1].width = Inches(4.5)
    tbl.columns[2].width = Inches(5.0)
    for r_idx, row in enumerate(rows):
        for c_idx, cell_text in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            cell.fill.solid()
            if r_idx == 0:
                cell.fill.fore_color.rgb = ACCENT
                color = WHITE
                bold = True
                size = 13
            else:
                cell.fill.fore_color.rgb = PANEL if r_idx % 2 == 0 else BG
                color = WHITE if c_idx == 2 else DIM
                bold = c_idx == 2
                size = 11
            tf = cell.text_frame
            tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
            tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            p.text = ""
            run = p.add_run()
            run.text = cell_text
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color


def slide_unique(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "獨有機制", "市面冇人咁做嘅 7 個設計",
               "唔係「我哋都有呢個 feature」、係 architecture 級別嘅唔同")
    items = [
        ("①", "Per-segment Baseline + Lazy Revert",
         "每段 zh_text 都有 baseline_zh 副本。Glossary 條目刪除 / 修改 → 下次 scan lazy revert。"),
        ("②", "Glossary Apply via LLM 智能替換",
         "唔係 string find-replace — LLM 看上下文保留語法時態，OpenCC s2twp 強制繁體輸出。"),
        ("③", "WYSIWYG Preview = Burn-in",
         "SVG paint-order + viewBox 1920×1080 + 共享 TTF + fontsdir 對齊 libass — pixel-perfect。"),
        ("④", "三段可組合 Pipeline (Profile = 樂高)",
         "ASR + 翻譯 + 字型 + Render 四模塊獨立 schema。動態參數面板，自由配對。"),
        ("⑤", "三層翻譯 Alignment",
         "Sentence pipeline + LLM-marker injection + opt-in Pass 2 enrichment — 業界冇人咁做。"),
        ("⑥", "結構化 QA Flags（schema 級別）",
         "[long] / [review] 喺 flags: List[str] 欄位、唔污染 zh_text 字串。永遠唔會錯燒入畫面。"),
        ("⑦", "Word-boundary + Smart-case Glossary Match",
         "「US」唔 match「must」/「trust」/「USA」。前後端 regex 完全 mirror — preview = backend 行為。"),
    ]
    y = Inches(2.0)
    col_h = Inches(0.65)
    for i, (num, title, body) in enumerate(items):
        ny = y + col_h * i
        # number bubble
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(0.7), ny + Inches(0.05),
                                     Inches(0.5), Inches(0.5))
        circle.fill.solid(); circle.fill.fore_color.rgb = ACCENT
        circle.line.fill.background(); circle.shadow.inherit = False
        tf = circle.text_frame
        tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num; r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = WHITE
        add_text(s, Inches(1.4), ny, Inches(4.5), Inches(0.4),
                 title, size=14, bold=True, color=ACCENT_2)
        add_text(s, Inches(5.9), ny + Inches(0.04), Inches(7.0), Inches(0.6),
                 body, size=11, color=WHITE)


def slide_techstack(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "Tech Stack", "全部開放、可審查、可遷移")
    cats = [
        ("Backend", [
            "Python 3.9+", "Flask + Flask-SocketIO",
            "faster-whisper / mlx-whisper", "Ollama + Qwen 3.5",
            "OpenRouter (Claude / GPT / Gemini)",
            "FFmpeg (ASS burn-in + MXF / XDCAM)",
            "OpenCC s2twp", "pySBD",
        ]),
        ("Frontend", [
            "Vanilla HTML / CSS / JS", "(冇 build step)",
            "Socket.IO client", "SVG paint-order overlay",
            "File System Access API", "Playwright async (E2E test)",
        ]),
        ("Data / Format", [
            "JSON file storage (plain, diff-able)",
            "CSV import / export (glossary)",
            "ASS / SRT / VTT / TXT subtitle output",
            "MP4 / MXF / MOV / MKV input",
            "MP4 H.264 / MXF ProRes / MXF XDCAM 422 output",
        ]),
    ]
    y = Inches(2.2)
    box_w = Inches(4.0); gap = Inches(0.2)
    for i, (cat, items) in enumerate(cats):
        x = Inches(0.6) + (box_w + gap) * i
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, box_w, Inches(0.07))
        bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background(); bar.shadow.inherit = False
        add_text(s, x, y + Inches(0.18), box_w, Inches(0.5),
                 cat, size=18, bold=True, color=WHITE)
        body = "\n".join(f"• {it}" for it in items)
        add_text(s, x, y + Inches(0.85), box_w, Inches(4),
                 body, size=12, color=DIM)


def slide_closing(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.6), SLIDE_W, Inches(2.5))
    panel.fill.solid(); panel.fill.fore_color.rgb = PANEL
    panel.line.fill.background(); panel.shadow.inherit = False
    add_accent_bar(s, Inches(0.6), Inches(2.85), Inches(0.12), Inches(2.0))
    add_text(s, Inches(0.9), Inches(2.85), Inches(11.5), Inches(0.5),
             "MoTitle", size=18, bold=True, color=ACCENT_2)
    add_text(s, Inches(0.9), Inches(3.35), Inches(12), Inches(1.4),
             "唔需要妥協，亦唔需要月費", size=48, bold=True, color=WHITE)
    add_text(s, Inches(0.9), Inches(4.55), Inches(12), Inches(0.5),
             "把控字幕生產嘅每一個 step、每一個 pixel、每一份數據",
             size=18, color=DIM)
    # Final value props
    chips_y = Inches(5.6)
    chips = [
        (" 100% 本地", SUCCESS, BG),
        (" 零月費", ACCENT, WHITE),
        (" 廣播交付級", ACCENT_2, BG),
        (" 模塊化 ", PANEL, WHITE),
        (" 開源 ", DIVIDER, WHITE),
    ]
    cx = Inches(0.9)
    for label, bg_c, fg_c in chips:
        add_chip(s, cx, chips_y, label, color=bg_c, fg=fg_c, size=12)
        cx += Inches(1.6)


# ============================================================
# Build
# ============================================================

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_funcs = [
        slide_cover,
        slide_problem,
        slide_solution,
        slide_pipeline,
        slide_dashboard,
        slide_profile_editor,
        slide_proofread_overview,
        slide_proofread_segment,
        slide_glossary_panel,
        slide_apply_modal,
        slide_revert,
        slide_subtitle_settings,
        slide_subtitle_overlay,
        slide_render_mp4,
        slide_render_mxf,
        slide_render_xdcam,
        slide_render_2pass,
        slide_comparison,
        slide_unique,
        slide_techstack,
        slide_closing,
    ]
    for fn in slide_funcs:
        fn(prs)

    # Add page footer to all slides except cover & closing
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        if i in (1, total):
            continue
        add_footer(slide, i, total)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"\nGenerated: {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
